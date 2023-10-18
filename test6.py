import pdfplumber
from pptx import Presentation
from pptx.util import Inches
from gensim.summarization.summarizer import summarize
#pip install gensim
def create_ppt_from_chapter_summaries(pdf_path, output_pptx_path):
    presentation = Presentation()
    chapter_content = []  # Store content for the current chapter
    i=0
    with pdfplumber.open(pdf_path) as pdf:
        for page_num in range(len(pdf.pages)):
            page = pdf.pages[page_num]
            page_text = page.extract_text()

            if page_text:
                # You can modify this condition to identify chapters based on text patterns
                if "Chapter" in page_text:
                    # Start of a new chapter
                    if chapter_content:
                        # Generate a summary for the previous chapter and add it to a slide
                        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide layout
                        title = slide.shapes.title
                        title.text = f"Chapter {len(presentation.slides)}"  # Use chapter number as the slide title

                        chapter_text = " ".join(chapter_content)  # Combine all pages in the chapter
                        summary = summarize(chapter_text)

                        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                        text_frame = text_box.text_frame
                        p = text_frame.add_paragraph()
                        p.text = summary

                    # Clear the content for the current chapter
                    chapter_content = []

                # Append the page text to the current chapter content
                chapter_content.append(page_text)

    # Generate a summary for the last chapter and add it to a slide
    if chapter_content:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide layout
        title = slide.shapes.title
        #title.text = f"Chapter {len(presentation.slides)}"  # Use chapter number as the slide title
        title.text = f"Slide {i}"  # Use i as title
        i+=i
        
        chapter_text = " ".join(chapter_content)  # Combine all pages in the last chapter
        summary = summarize(chapter_text)

        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = summary

    presentation.save(output_pptx_path)

# Set the path to your PDF file and the output PPTX file
pdf_path = "ML2.pdf"
output_pptx_path = "output.pptx"
create_ppt_from_chapter_summaries(pdf_path, output_pptx_path)
