import PyPDF2
from pptx import Presentation
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.probability import FreqDist
from heapq import nlargest
from transformers import pipeline

# extract text from pdf
def extract_text_from_pdf(pdf_file_path):
    text = ""
    try:
        with open(pdf_file_path, 'rb') as pdf_file:
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += page.extract_text()
        return text
    except Exception as e:
        print(f"Error reading PDF: {str(e)}")
        return None

# Function to generate a PowerPoint presentation (Placeholder)
def generate_ppt_from_text(text_content):
    presentation = Presentation()
    for section in text_content.split('\n\n'):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = section[:100]  # Use the first 100 characters as slide title
        slide.shapes.add_textbox(left=100, top=100, width=500, height=300).text = section
    presentation.save('output.pptx')

# Function to integrate YouTube video links (Placeholder)
def integrate_youtube_videos(text_content):
    # Placeholder for YouTube integration code
    pass

def generate_summary(text_content, percentage=0.8):
    # Tokenize the text into sentences and words
    sentences = sent_tokenize(text_content)
    words = word_tokenize(text_content)
    
    # Remove stopwords and punctuation
    stop_words = set(stopwords.words('english'))
    words = [word.lower() for word in words if word.isalnum() and word.lower() not in stop_words]
    
    # Calculate word frequencies
    word_freq = FreqDist(words)
    
    # Calculate the score for each sentence based on word frequency
    sentence_scores = {}
    for sentence in sentences:
        for word in word_tokenize(sentence):
            if word.lower() in word_freq:
                if sentence not in sentence_scores:
                    sentence_scores[sentence] = word_freq[word.lower()]
                else:
                    sentence_scores[sentence] += word_freq[word.lower()]
    
    # Calculate the desired number of sentences for the summary based on the percentage
    num_sentences = int(len(sentences) * percentage)
    
    # Get the top 'num_sentences' sentences with the highest scores
    summary_sentences = nlargest(num_sentences, sentence_scores, key=sentence_scores.get)
    
    # Combine the selected sentences to generate the summary
    summary = ' '.join(summary_sentences)
    summarizer = pipeline("summarization")
    summary = summarizer(text_content, max_length=150, min_length=30, do_sample=False)[0]['summary_text']
    return summary

# Example usage
pdf_file_path = 'eeen101.pdf'
text_content = extract_text_from_pdf(pdf_file_path)

if text_content:
    generate_ppt_from_text(text_content)
    integrate_youtube_videos(text_content)
    summary = generate_summary(text_content)
    print(summary)
