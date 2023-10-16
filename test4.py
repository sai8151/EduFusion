import spacy
import pdfplumber
from pptx import Presentation
from pptx.util import Inches
# Initialize SpaCy with TextRank
nlp = spacy.load("en_core_web_lg")

def extract_headings_and_content(pdf_path):
    headings_and_content = []
    with pdfplumber.open(pdf_path) as pdf:
        for page_num in range(len(pdf.pages)):
            page = pdf.pages[page_num]
            page_text = page.extract_text()
            if page_text:
                # Split the page text into lines
                lines = page_text.split('\n')
                headings = []
                content = []
                for line in lines:
                    # You can modify this condition to identify headings based on text patterns
                    if line.isupper() or line.startswith('Chapter') or line.startswith('Section'):
                        headings.append(line)
                    else:
                        content.append(line)
                if headings:
                    headings_and_content.append((headings, content))
    return headings_and_content

def create_ppt_from_headings(pdf_path, output_pptx_path):
    headings_and_content = extract_headings_and_content(pdf_path)
    presentation = Presentation()
    for headings, content in headings_and_content:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide layout
        if headings:
            title = slide.shapes.title
            title.text = headings[0]  # Use the first heading as the slide title
        if content:
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
            text_frame = text_box.text_frame
            for line in content:
                p = text_frame.add_paragraph()
                p.text = line
    presentation.save(output_pptx_path)

# Set the path to your PDF file and the output PPTX file
pdf_path = "ML2.pdf"
output_pptx_path = "output.pptx"
create_ppt_from_headings(pdf_path, output_pptx_path)
