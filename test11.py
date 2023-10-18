import pdfplumber
from pptx import Presentation
from pptx.util import Inches
import spacy
from spacy.lang.en import English
from spacy.lang.en.stop_words import STOP_WORDS
from pptx.util import Inches, Pt 
# Load the spaCy language model
nlp = spacy.load("en_core_web_lg")
stop_words = STOP_WORDS

# Function to summarize text into 3 to 5 bullet points
def summarize_text(text):
    doc = nlp(text)
    sentences = [sent.text for sent in doc.sents]

    # Calculate the term frequency for each word
    word_freq = {}
    for word in doc:
        if word.text.lower() not in stop_words:
            if word.text not in word_freq:
                word_freq[word.text] = 1
            else:
                word_freq[word.text] += 1

    # Calculate the sentence score based on word frequency
    sentence_scores = {}
    for sent in sentences:
        for word in nlp(sent):
            if word.text.lower() in word_freq:
                if sent not in sentence_scores:
                    sentence_scores[sent] = word_freq[word.text.lower()]
                else:
                    sentence_scores[sent] += word_freq[word.text.lower()]

    # Get the top 3 to 5 sentences with the highest scores
    summary_sentences = sorted(sentence_scores, key=sentence_scores.get, reverse=True)[:5]
    return summary_sentences

def create_ppt_from_chapter_summaries(pdf_path, output_pptx_path):
    presentation = Presentation()
    chapter_content = []  # Store content for the current chapter
    subtopic_content = []  # Store content for the current subtopic
    chapter_title = None
    subtopic_title = None
    i = 1  # Initialize the slide number

    with pdfplumber.open(pdf_path) as pdf:
        for page_num in range(len(pdf.pages)):
            page = pdf.pages[page_num]
            page_text = page.extract_text()

            if page_text:
                # Check if it's a chapter title
                if "Chapter" in page_text:
                    # Start of a new chapter
                    if chapter_content:
                        # Generate a summary for the chapter and add it to a slide
                        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide layout
                        title = slide.shapes.title
                        title.text = chapter_title  # Use chapter title as the slide title
                        title.text_frame.paragraphs[0].font.size = Pt(44)  # Set title font size to 44pt

                        chapter_text = " ".join(chapter_content)  # Combine all pages in the chapter
                        summary_sentences = summarize_text(chapter_text)

                        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
                        text_frame = text_box.text_frame
                        for sentence in summary_sentences:
                            p = text_frame.add_paragraph()
                            p.text = sentence

                    # Clear the content for the current chapter and subtopic
                    chapter_content = []
                    subtopic_content = []
                    chapter_title = page_text.strip()

                # Check if it's a subtopic
                elif page_text.lower().startswith("section") or page_text.lower().startswith("subtopic"):
                    # Start of a new subtopic
                    if subtopic_content:
                        # Generate a summary for the subtopic and add it to a slide
                        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide layout
                        title = slide.shapes.title
                        title.text = subtopic_title  # Use subtopic title as the slide title
                        title.text_frame.paragraphs[0].font.size = Pt(44)  # Set title font size to 44pt

                        subtopic_text = " ".join(subtopic_content)  # Combine all pages in the subtopic
                        summary_sentences = summarize_text(subtopic_text)

                        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
                        text_frame = text_box.text_frame
                        for sentence in summary_sentences:
                            p = text_frame.add_paragraph()
                            p.text = sentence

                    # Clear the content for the current subtopic
                    subtopic_content = []
                    subtopic_title = page_text.strip()

                # Append the page text to the current chapter or subtopic content
                if chapter_title:
                    chapter_content.append(page_text)
                if subtopic_title:
                    subtopic_content.append(page_text)

    # Generate a summary for the last chapter and add it to a slide
    if chapter_content:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide layout
        title = slide.shapes.title
        title.text = chapter_title  # Use chapter title as the slide title
        title.text_frame.paragraphs[0].font.size = Pt(44)  # Set title font size to 44pt

        chapter_text = " ".join(chapter_content)  # Combine all pages in the last chapter
        summary_sentences = summarize_text(chapter_text)

        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        text_frame = text_box.text_frame
        for sentence in summary_sentences:
            p = text_frame.add_paragraph()
            p.text = sentence

    presentation.save(output_pptx_path)

# Set the path to your PDF file and the output PPTX file
pdf_path = "ML2.pdf"
output_pptx_path = "output.pptx"
create_ppt_from_chapter_summaries(pdf_path, output_pptx_path)
